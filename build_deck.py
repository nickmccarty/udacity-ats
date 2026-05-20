"""
build_deck.py -- SQLite Search Cache micro-lesson PPTX

Formatting aligned with "Copy of Udacity Teaching Sample Slide Template 10.23.pptx":
  - Slide size: 20" x 11.25"  (template dimensions, 16:9)
  - Font: Open Sans (body), Roboto Mono (code)
  - Colors: Udacity light theme -- lt1 (#FFFFFF) bg, dk2 (#171A53) headings
  - Layout grid: header band y=0-1.3", content body y=2.0"+

11 slides mirroring the 11 video segments exactly.

Run:
    python micro-lesson/course-4/build_deck.py
Output:
    micro-lesson/course-4/sqlite-search-cache.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Palette  -- Udacity template light theme (lt1 background)
# ---------------------------------------------------------------------------
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)   # lt1 -- slide background
LGRAY    = RGBColor(0xF6, 0xF6, 0xF6)  # lt2 -- card surface
MGRAY    = RGBColor(0xF0, 0xF0, 0xF0)  # mid-gray -- alternate card
DARK     = RGBColor(0x0B, 0x0B, 0x0B)  # dk1 -- body text
NAVY     = RGBColor(0x17, 0x1A, 0x53)  # dk2 -- headings, header band
BLUE     = RGBColor(0x20, 0x15, 0xFF)  # accent1 -- primary CTA
TEAL     = RGBColor(0x00, 0xC5, 0xA1)  # accent2
LIME     = RGBColor(0xBD, 0xEA, 0x09)  # accent4
LBLUE    = RGBColor(0x65, 0x97, 0xFF)  # accent5
LPUR     = RGBColor(0xB1, 0x81, 0xFF)  # accent6

# Derived surfaces
BG       = WHITE
BG_CARD  = LGRAY                         # F6F6F6 -- card fill
BG_CODE  = RGBColor(0xEA, 0xEE, 0xF2)   # light blue-gray code block bg
BORDER   = RGBColor(0xDB, 0xE2, 0xE8)   # accent3 -- card border
MUTED    = RGBColor(0x44, 0x44, 0x66)   # secondary body text
DIM      = RGBColor(0x88, 0x88, 0xAA)   # labels / tertiary
CODE_CLR = NAVY                          # code text on light bg

RED      = RGBColor(0xDC, 0x26, 0x26)
RED_BG   = RGBColor(0xFE, 0xE2, 0xE2)
RED_FG   = RGBColor(0x99, 0x1B, 0x1B)
TEAL_BG  = RGBColor(0xE6, 0xFA, 0xF7)
TEAL_FG  = RGBColor(0x00, 0x7A, 0x66)
BLUE_BG  = RGBColor(0xEE, 0xF2, 0xFF)
BLUE_FG  = RGBColor(0x1D, 0x1A, 0xCC)
LIME_BG  = RGBColor(0xF7, 0xFA, 0xD6)
LIME_FG  = RGBColor(0x4A, 0x58, 0x00)

FONT      = "Open Sans"
FONT_CODE = "Roboto Mono"

# ---------------------------------------------------------------------------
# Slide canvas  -- matches template dimensions exactly
# ---------------------------------------------------------------------------
W = Inches(20)      # 20.000" -- template width
H = Inches(11.25)   # 11.250" -- template height

MARGIN_L   = Inches(0.75)
MARGIN_R   = Inches(0.75)
CONTENT_W  = W - MARGIN_L - MARGIN_R    # 18.5"
HEADER_TOP = Inches(0.18)               # chip y-position
TITLE_TOP  = Inches(1.35)               # slide title y-position
BODY_TOP   = Inches(2.2)                # main body content y-position
BODY_H     = H - BODY_TOP - Inches(0.5)


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(sl, color=BG):
    s = sl.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def box(sl, x, y, w, h, fill, border=None, bpt=0.75):
    s = sl.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bpt)
    else:
        s.line.fill.background()
    return s

def accent_bar(sl, color=NAVY):
    """Top accent bar -- navy by default (matches template header band)."""
    box(sl, 0, 0, W, Inches(0.1), color)

def header_band(sl):
    """Light gray header band behind chip area."""
    box(sl, 0, 0, W, Inches(0.75), LGRAY)
    accent_bar(sl)

def chip_label(sl, text, color=NAVY):
    """Section chip at fixed header position."""
    cw = Inches(3.2)
    ch = Inches(0.38)
    cx = MARGIN_L
    cy = HEADER_TOP + Inches(0.18)
    box(sl, cx, cy, cw, ch, color)
    t = sl.shapes.add_textbox(cx + Inches(0.14), cy + Inches(0.04),
                               cw - Inches(0.2), ch - Inches(0.06))
    p = t.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text; r.font.name = FONT
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE

def slide_num_label(sl, num):
    txt(sl, W - Inches(2.5), HEADER_TOP + Inches(0.24),
        Inches(2.3), Inches(0.28),
        f"Slide {num} of 11", 12, color=DIM, align=PP_ALIGN.RIGHT)

def txt(sl, x, y, w, h, text, size=24, bold=False, italic=False,
        color=DARK, font=FONT, align=PP_ALIGN.LEFT, wrap=True):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.name = font
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def multi_para(sl, x, y, w, h, lines, size=24, color=DARK, font=FONT,
               bold_first=False):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line; r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold_first and i == 0
        r.font.color.rgb = color
    return tb

def code_block(sl, x, y, w, h, text, size=15):
    box(sl, x, y, w, h, BG_CODE, border=BORDER, bpt=0.6)
    box(sl, x, y, Inches(0.045), h, NAVY)
    tb = sl.shapes.add_textbox(x + Inches(0.25), y + Inches(0.2),
                                w - Inches(0.35), h - Inches(0.35))
    tb.word_wrap = False
    tf = tb.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text; r.font.name = FONT_CODE
    r.font.size = Pt(size); r.font.color.rgb = CODE_CLR

def card(sl, x, y, w, h, fill=BG_CARD, accent_col=None):
    box(sl, x, y, w, h, fill, border=BORDER, bpt=0.5)
    if accent_col:
        box(sl, x, y, w, Inches(0.055), accent_col)

def callout(sl, x, y, w, h, text, style="blue", size=20):
    cfg = {
        "blue": (BLUE_BG, BLUE,  BLUE_FG),
        "teal": (TEAL_BG, TEAL,  TEAL_FG),
        "red":  (RED_BG,  RED,   RED_FG),
        "lime": (LIME_BG, LIME,  LIME_FG),
        "navy": (LGRAY,   NAVY,  NAVY),
    }
    bg_, accent, fg = cfg.get(style, cfg["blue"])
    box(sl, x, y, w, h, bg_, border=BORDER, bpt=0.4)
    box(sl, x, y, Inches(0.055), h, accent)
    txt(sl, x + Inches(0.2), y + Inches(0.15),
        w - Inches(0.3), h - Inches(0.25),
        text, size=size, color=fg, wrap=True)

def notes(sl, text):
    sl.notes_slide.notes_text_frame.text = text


# ===========================================================================
# Slide builders  (11 slides matching video segments)
# ===========================================================================

# ---------------------------------------------------------------------------
# Slide 1  --  The Problem  [seg 1, 22s]
# ---------------------------------------------------------------------------
def slide_01(prs):
    sl = blank(prs)
    fill_bg(sl)
    header_band(sl)
    accent_bar(sl, RED)
    chip_label(sl, "THE PROBLEM", RED)
    slide_num_label(sl, 1)

    txt(sl, MARGIN_L, TITLE_TOP, CONTENT_W, Inches(1.5),
        "A cost & latency problem\ninvisible in single-session work",
        size=48, bold=True, color=NAVY)

    txt(sl, MARGIN_L, Inches(3.1), Inches(13.5), Inches(0.9),
        "Multi-session auto-research re-issues the same web search queries on every "
        "loop -- burning API credits, hitting rate limits, adding latency per session "
        "with identical round trips.",
        size=24, color=MUTED)

    labels = ["COST", "RATE LIMITS", "LATENCY"]
    bodies = [
        "API credits consumed for results already retrieved in session 1",
        "Same queries push you toward provider throttle thresholds",
        "30-90 s of round-trip time added per session -- for zero new information",
    ]
    for i in range(3):
        cx = MARGIN_L + Inches(i * 6.2)
        card(sl, cx, Inches(4.3), Inches(5.9), Inches(5.9), accent_col=RED)
        txt(sl, cx + Inches(0.3), Inches(4.6), Inches(5.3), Inches(0.4),
            labels[i], 13, bold=True, color=RED)
        txt(sl, cx + Inches(0.3), Inches(5.2), Inches(5.3), Inches(4.5),
            bodies[i], 22, color=MUTED)

    notes(sl, (
        "This lesson solves a cost and latency problem invisible in single-session "
        "exercises, but dominant in production auto-research workflows."
    ))
    return sl


# ---------------------------------------------------------------------------
# Slide 2  --  The Solution  [seg 2, 29s]
# ---------------------------------------------------------------------------
def slide_02(prs):
    sl = blank(prs)
    fill_bg(sl)
    header_band(sl)
    chip_label(sl, "THE SOLUTION", NAVY)
    slide_num_label(sl, 2)

    txt(sl, MARGIN_L, TITLE_TOP, CONTENT_W, Inches(1.0),
        "The fix: a SHA-256 keyed SQLite cache",
        size=44, bold=True, color=NAVY)

    card(sl, MARGIN_L, BODY_TOP, Inches(7.8), Inches(4.5))
    txt(sl, MARGIN_L + Inches(0.3), BODY_TOP + Inches(0.3), Inches(7.2), Inches(0.5),
        "EVERY QUERY", 14, bold=True, color=DIM)
    txt(sl, MARGIN_L + Inches(0.3), BODY_TOP + Inches(1.0), Inches(7.2), Inches(0.9),
        "search(q)", 32, bold=True, color=NAVY, font=FONT_CODE)
    txt(sl, MARGIN_L + Inches(0.3), BODY_TOP + Inches(2.1), Inches(7.2), Inches(0.5),
        "session 1  .  2  .  3  .  N", 18, color=DIM)

    txt(sl, Inches(8.8), BODY_TOP + Inches(1.6), Inches(1.0), Inches(0.8),
        "->", 34, bold=True, color=DIM, font=FONT_CODE)

    card(sl, Inches(10.0), BODY_TOP, Inches(9.75), Inches(4.5), accent_col=NAVY)
    txt(sl, Inches(10.3), BODY_TOP + Inches(0.3), Inches(9.1), Inches(0.5),
        "SQLITE CACHE", 14, bold=True, color=NAVY)
    callout(sl, Inches(10.3), BODY_TOP + Inches(1.0), Inches(4.4), Inches(1.35),
            "HIT: return stored result\n0 API calls", "teal", 20)
    callout(sl, Inches(15.0), BODY_TOP + Inches(1.0), Inches(4.5), Inches(1.35),
            "MISS: call API -> store -> return", "red", 20)

    txt(sl, MARGIN_L, Inches(7.2), Inches(3.0), Inches(2.2),
        "0", 110, bold=True, color=TEAL)
    txt(sl, Inches(3.2), Inches(8.1), Inches(5.5), Inches(0.7),
        "API calls -- sessions 2 through N", 22, color=MUTED)

    badge_data = [("SHA-256 key", NAVY), ("SQLite storage", NAVY),
                  ("TTL eviction", TEAL), ("Provider-agnostic", DIM)]
    for i, (label, col) in enumerate(badge_data):
        bx = Inches(9.5 + i * 2.65)
        box(sl, bx, Inches(8.15), Inches(2.5), Inches(0.5), BG_CARD, border=col, bpt=0.5)
        txt(sl, bx + Inches(0.15), Inches(8.22), Inches(2.2), Inches(0.38),
            label, 15, color=col)

    notes(sl, "SHA-256 keyed SQLite cache; sessions 2-N make 0 API calls.")
    return sl


# ---------------------------------------------------------------------------
# Slide 3  --  Try It First  [seg 3, 28s]
# ---------------------------------------------------------------------------
def slide_03(prs):
    sl = blank(prs)
    fill_bg(sl)
    header_band(sl)
    accent_bar(sl, RED)
    chip_label(sl, "TRY-AND-FAIL", RED)
    slide_num_label(sl, 3)

    txt(sl, MARGIN_L, TITLE_TOP, Inches(11), Inches(0.9),
        "Try it first -- before the explanation",
        size=36, bold=True, color=NAVY)

    code_block(sl, MARGIN_L, BODY_TOP, Inches(11.5), Inches(7.5), (
        "api_calls = 0\n\n"
        "def search(query, max_results=10):\n"
        "    global api_calls\n"
        "    api_calls += 1\n"
        "    print(f'[API #{api_calls}] {query[:55]}')\n"
        "    return list(DDGS().text(query, max_results=max_results))\n\n"
        "QUERIES = [\n"
        "    'transformer attention mechanisms survey',\n"
        "    'KV cache optimization techniques',\n"
        "    'flash attention implementation pytorch',\n"
        "    'multi-head latent attention MLA',\n"
        "]\n\n"
        "for session in range(1, 4):   # 3 sessions\n"
        "    for q in QUERIES:\n"
        "        search(q)\n\n"
        "print(f'Total: {api_calls} calls')"
    ), size=17)

    callout(sl, Inches(12.55), BODY_TOP, Inches(7.2), Inches(2.5),
            "Three auto-research sessions.\n"
            "Same four queries each session.\n"
            "Run it -- count the API calls.", "navy", 22)

    txt(sl, Inches(12.55), Inches(5.2), Inches(7.2), Inches(2.0),
        "How many of those calls carry new information?",
        34, bold=True, color=NAVY)

    notes(sl, "Three sessions, four queries each. Run before explanation.")
    return sl


# ---------------------------------------------------------------------------
# Slide 4  --  The Waste  [seg 4, 45s]
# ---------------------------------------------------------------------------
def slide_04(prs):
    sl = blank(prs)
    fill_bg(sl)
    header_band(sl)
    accent_bar(sl, RED)
    chip_label(sl, "TRY-AND-FAIL", RED)
    slide_num_label(sl, 4)

    txt(sl, MARGIN_L, TITLE_TOP, Inches(5.0), Inches(0.5),
        "The result", 20, bold=True, color=DIM)

    txt(sl, MARGIN_L, Inches(2.0), Inches(4.5), Inches(3.5),
        "12", 130, bold=True, color=RED)
    txt(sl, MARGIN_L, Inches(5.8), Inches(4.5), Inches(0.55),
        "total API calls", 18, color=DIM)

    card(sl, MARGIN_L, Inches(6.55), Inches(8.2), Inches(3.5))
    multi_para(sl, MARGIN_L + Inches(0.3), Inches(6.85), Inches(7.6), Inches(2.9), [
        "4 queries x 3 sessions = 12 calls",
        "Novel results -- session 1:      4",
        "Novel results -- sessions 2-3:   0",
        "Redundant calls:  8  (67%)",
    ], size=22, color=MUTED, font=FONT_CODE)

    callout(sl, MARGIN_L, Inches(10.2), Inches(8.2), Inches(0.85),
            "No exception. No warning. No log entry. The waste is completely silent.", "red", 20)

    txt(sl, Inches(9.7), TITLE_TOP, Inches(9.5), Inches(0.5),
        "At production scale", 20, bold=True, color=MUTED)
    card(sl, Inches(9.7), Inches(2.0), Inches(9.5), Inches(4.5))
    multi_para(sl, Inches(10.0), Inches(2.35), Inches(8.9), Inches(3.7), [
        "12 queries x 10 sessions",
        "= 108 redundant calls",
        "= 30-90 s added latency",
        "  every time auto-research runs",
    ], size=24, color=MUTED, font=FONT_CODE)

    callout(sl, Inches(9.7), Inches(6.7), Inches(9.5), Inches(1.9),
            "The cache eliminates this completely.\n"
            "Sessions 2-N cost 0 API calls for any query already seen.", "teal", 22)

    notes(sl, "12 calls, only 4 novel. 108 redundant at scale.")
    return sl


# ---------------------------------------------------------------------------
# Slide 5  --  SHA-256 Sensitivity  [seg 5, 30s]
# ---------------------------------------------------------------------------
def slide_05(prs):
    sl = blank(prs)
    fill_bg(sl)
    header_band(sl)
    chip_label(sl, "CONCEPT: NORMALISATION", NAVY)
    slide_num_label(sl, 5)

    txt(sl, MARGIN_L, TITLE_TOP, CONTENT_W, Inches(1.0),
        "Core insight: SHA-256 is deterministic -- and case-sensitive",
        size=38, bold=True, color=NAVY)

    callout(sl, MARGIN_L, BODY_TOP, CONTENT_W, Inches(1.6),
            "Same bytes -> same digest. Always. But a single uppercase letter or an "
            "extra space produces a completely different 256-bit output.", "navy", 24)

    code_block(sl, MARGIN_L, Inches(4.15), CONTENT_W, Inches(2.8), (
        "# Without normalisation -- three different cache keys:\n"
        "sha256('Flash Attention')    = 'a1b2c3...'   # session 1\n"
        "sha256('flash attention')    = '9f8e7d...'   # different!\n"
        "sha256('flash  attention')   = '4d5e6f...'   # extra space = third key"
    ), size=18)

    callout(sl, MARGIN_L, Inches(7.2), CONTENT_W, Inches(1.6),
            "Semantically identical queries -> different cache keys -> cache always misses. "
            "Every session makes the full API call.", "red", 24)

    notes(sl, "SHA-256 is case- and whitespace-sensitive. Normalise before hashing.")
    return sl


# ---------------------------------------------------------------------------
# Slide 6  --  Normalisation Pipeline  [seg 6, 45s]
# ---------------------------------------------------------------------------
def slide_06(prs):
    sl = blank(prs)
    fill_bg(sl)
    header_band(sl)
    chip_label(sl, "CONCEPT: NORMALISATION", NAVY)
    slide_num_label(sl, 6)

    txt(sl, MARGIN_L, TITLE_TOP, CONTENT_W, Inches(0.9),
        "The normalisation pipeline -- three steps",
        size=38, bold=True, color=NAVY)

    steps = [
        (".lower()",  "STEP 1",             "Folds ALL uppercase to lowercase.\n'Flash' -> 'flash'",    BLUE),
        (".split()",  "STEP 2 -- KEY STEP", "Splits on ANY whitespace: tabs, double spaces, newlines.\nNo argument = all whitespace.", TEAL),
        (" .join()",  "STEP 3",             "Rejoins with exactly one space.\nCanonical form guaranteed.", BLUE),
    ]
    for i, (fn, label, body, col) in enumerate(steps):
        sx = MARGIN_L + Inches(i * 6.35)
        card(sl, sx, BODY_TOP, Inches(6.0), Inches(3.9), accent_col=col)
        txt(sl, sx + Inches(0.3), BODY_TOP + Inches(0.3), Inches(5.4), Inches(0.4),
            label, 13, bold=True, color=col)
        txt(sl, sx + Inches(0.3), BODY_TOP + Inches(0.9), Inches(5.4), Inches(0.75),
            fn, 28, bold=True, color=NAVY, font=FONT_CODE)
        txt(sl, sx + Inches(0.3), BODY_TOP + Inches(1.9), Inches(5.4), Inches(1.8),
            body, 20, color=MUTED)

    code_block(sl, MARGIN_L, Inches(6.5), CONTENT_W, Inches(1.85), (
        "def _cache_key(query: str) -> str:\n"
        "    normalised = ' '.join(query.lower().split())\n"
        "    return hashlib.sha256(normalised.encode()).hexdigest()"
    ), size=19)

    callout(sl, MARGIN_L, Inches(8.55), CONTENT_W, Inches(1.35),
            "Rule: every cache key derived from user input should normalise before hashing. "
            "Applies to databases, rate-limiters, and dedup systems alike.", "teal", 22)

    notes(sl, ".lower() -> .split() -> ' '.join() -> sha256()")
    return sl


# ---------------------------------------------------------------------------
# Slide 7  --  cached_search() Logic  [seg 7, 35s]
# ---------------------------------------------------------------------------
def slide_07(prs):
    sl = blank(prs)
    fill_bg(sl)
    header_band(sl)
    chip_label(sl, "DEMO: cached_search()", NAVY)
    slide_num_label(sl, 7)

    txt(sl, MARGIN_L, TITLE_TOP, Inches(11.5), Inches(0.9),
        "The public interface -- four lines of logic",
        size=36, bold=True, color=NAVY)

    code_block(sl, MARGIN_L, BODY_TOP, Inches(11.5), Inches(7.65), (
        "def cached_search(\n"
        "    query: str,\n"
        "    search_fn: Callable[[str, int], list[dict]],\n"
        "    ttl: int = DEFAULT_TTL,\n"
        "    max_results: int = 10,\n"
        ") -> list[dict]:\n"
        "\n"
        "    # HIT path\n"
        "    cached = get(query)\n"
        "    if cached is not None:\n"
        "        print(f'[cache HIT ] {query[:55]}')\n"
        "        return cached\n"
        "\n"
        "    # MISS path\n"
        "    print(f'[cache MISS] {query[:55]}')\n"
        "    results = search_fn(query, max_results)\n"
        "    if results:               # never cache empty\n"
        "        put(query, results, ttl=ttl)\n"
        "    return results"
    ), size=17)

    callout(sl, Inches(12.55), BODY_TOP, Inches(7.2), Inches(2.6),
            "HIT path:\nget() returns stored results.\nNo network call. Instant return.", "teal", 22)

    callout(sl, Inches(12.55), Inches(5.3), Inches(7.2), Inches(2.6),
            "MISS path:\nCall the search provider.\nStore in cache. Return results.", "red", 22)

    card(sl, Inches(12.55), Inches(8.2), Inches(7.2), Inches(1.65))
    txt(sl, Inches(12.85), Inches(8.4), Inches(6.6), Inches(0.35),
        "ONE-LINE CALLSITE CHANGE", 14, bold=True, color=DIM)
    txt(sl, Inches(12.85), Inches(8.85), Inches(6.6), Inches(0.8),
        "search(q)  ->  cached_search(q, search_fn=search)",
        17, color=NAVY, font=FONT_CODE)

    notes(sl, "Four lines: get() hit check, print HIT/MISS, search_fn call, put().")
    return sl


# ---------------------------------------------------------------------------
# Slide 8  --  Design Choices  [seg 8, 45s]
# ---------------------------------------------------------------------------
def slide_08(prs):
    sl = blank(prs)
    fill_bg(sl)
    header_band(sl)
    chip_label(sl, "DEMO: cached_search()", NAVY)
    slide_num_label(sl, 8)

    txt(sl, MARGIN_L, TITLE_TOP, CONTENT_W, Inches(0.9),
        "4 design choices in cached_search()",
        size=40, bold=True, color=NAVY)

    choices = [
        ("1 -- search_fn: Callable",  BLUE,
         "Provider-agnostic. Pass DDGS, Tavily, Bing, or a test stub.\n"
         "The cache doesn't care which search backend you use."),
        ("2 -- Never cache empty",     TEAL,
         "if results: put(...)  An empty list usually means the API failed "
         "or returned nothing. Caching it blocks retries for the full TTL window."),
        ("3 -- TTL per query",         BLUE,
         "Short TTL for breaking news queries.\nLong TTL for evergreen research topics.\n"
         "Default: 86400 s (24 h)."),
        ("4 -- Structured HIT/MISS",   TEAL,
         "[cache HIT ] / [cache MISS]  Structured labels that survive log aggregation. "
         "Easy to grep across multi-session logs to measure hit rate."),
    ]
    positions = [
        (MARGIN_L,      BODY_TOP),
        (Inches(10.4),  BODY_TOP),
        (MARGIN_L,      Inches(6.8)),
        (Inches(10.4),  Inches(6.8)),
    ]
    for (label, col, body), (cx, cy) in zip(choices, positions):
        card(sl, cx, cy, Inches(9.1), Inches(3.9), accent_col=col)
        txt(sl, cx + Inches(0.3), cy + Inches(0.3), Inches(8.5), Inches(0.45),
            label, 14, bold=True, color=col)
        txt(sl, cx + Inches(0.3), cy + Inches(0.95), Inches(8.5), Inches(2.65),
            body, 21, color=MUTED)

    notes(sl, "Callable, no empty cache, TTL per type, structured labels.")
    return sl


# ---------------------------------------------------------------------------
# Slide 9  --  Quiz  [seg 9, 70s]
# ---------------------------------------------------------------------------
def slide_09(prs):
    sl = blank(prs)
    fill_bg(sl)
    header_band(sl)
    chip_label(sl, "QUIZ", NAVY)
    slide_num_label(sl, 9)

    card(sl, MARGIN_L, TITLE_TOP, CONTENT_W, Inches(2.2))
    txt(sl, MARGIN_L + Inches(0.3), TITLE_TOP + Inches(0.2),
        CONTENT_W - Inches(0.6), Inches(1.8),
        "Your cached_search() uses SHA-256 of the raw query string. "
        "A user submits \"Flash Attention\" in session 1 and \"flash attention\" in session 2. "
        "The cache returns MISS and makes a second API call. "
        "What is the root cause and the correct fix?",
        22, bold=True, color=NAVY)

    options = [
        ("A", "Use SQLite LIKE instead of WHERE key=? so the lookup is case-insensitive.", False),
        ("B", "Key computed from the raw string. Fix: normalise with ' '.join(query.lower().split()) before hashing.", True),
        ("C", "Increase the TTL so the second session is more likely to find a live entry.", False),
        ("D", "Store the raw query string as the key instead of a hash.", False),
    ]
    for i, (letter, text, correct) in enumerate(options):
        oy = Inches(3.8 + i * 1.7)
        accent = TEAL if correct else BORDER
        fill   = TEAL_BG if correct else BG_CARD
        box(sl, MARGIN_L, oy, CONTENT_W, Inches(1.5), fill, border=accent, bpt=0.75)
        box(sl, MARGIN_L, oy, Inches(0.6), Inches(1.5), accent)
        txt(sl, MARGIN_L + Inches(0.08), oy + Inches(0.35), Inches(0.5), Inches(0.75),
            letter, 22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, MARGIN_L + Inches(0.8), oy + Inches(0.25),
            CONTENT_W - Inches(1.0), Inches(0.95),
            text, 20, color=TEAL_FG if correct else MUTED)

    notes(sl, (
        "B is correct. Normalise with ' '.join(query.lower().split()) before hashing.\n"
        "A: LIKE on a hash column matches nothing. C: TTL != key identity. D: raw strings don't fix case."
    ))
    return sl


# ---------------------------------------------------------------------------
# Slide 10  --  Recap  [seg 10, 65s]
# ---------------------------------------------------------------------------
def slide_10(prs):
    sl = blank(prs)
    fill_bg(sl)
    header_band(sl)
    accent_bar(sl, TEAL)
    chip_label(sl, "RECAP", TEAL)
    slide_num_label(sl, 10)

    txt(sl, MARGIN_L, TITLE_TOP, CONTENT_W, Inches(1.0),
        "What you can do now that you couldn't before:",
        size=38, bold=True, color=NAVY)

    points = [
        ("Explain",   "Why repeated auto-research queries waste API calls across sessions "
                      "and why the failure is silent."),
        ("Normalise", "Queries with ' '.join(query.lower().split()) before SHA-256 -- "
                      "same intent always maps to the same cache key."),
        ("Implement", "get() with expiry check and put() with upsert + inline eviction."),
        ("Wrap",      "Any search callable with cached_search(q, search_fn) in a one-line "
                      "change -- provider-agnostic and TTL-tunable."),
    ]
    for i, (verb, body) in enumerate(points):
        py = BODY_TOP + Inches(i * 1.7)
        box(sl, MARGIN_L, py + Inches(0.45), Inches(0.12), Inches(0.12), TEAL)
        txt(sl, MARGIN_L + Inches(0.3), py, Inches(2.8), Inches(0.55),
            verb + ":", 24, bold=True, color=NAVY)
        txt(sl, MARGIN_L + Inches(0.3), py + Inches(0.6),
            CONTENT_W - Inches(0.4), Inches(0.95),
            body, 22, color=MUTED)

    callout(sl, MARGIN_L, Inches(9.3), CONTENT_W, Inches(1.55),
            "Sessions 2-N: 0 API calls for any repeated query.  "
            "A 12-query loop x 10 sessions = 108 API calls saved per topic per day.", "teal", 23)

    notes(sl, "Explain, normalise, implement, wrap -- all in one-line change.")
    return sl


# ---------------------------------------------------------------------------
# Slide 11  --  Outro  [seg 11, 14.22s]
# ---------------------------------------------------------------------------
def slide_11(prs):
    sl = blank(prs)
    fill_bg(sl)
    header_band(sl)
    accent_bar(sl, TEAL)

    txt(sl, 0, Inches(3.6), W, Inches(1.5),
        "Thank you for watching.", 52, bold=True, color=NAVY,
        align=PP_ALIGN.CENTER)
    txt(sl, 0, Inches(5.4), W, Inches(0.85),
        "Lesson materials -- starter code & solution -- in micro-lesson/course-4/",
        26, color=MUTED, align=PP_ALIGN.CENTER)
    txt(sl, 0, Inches(6.5), W, Inches(1.0),
        "See you in the next lesson.", 34, bold=True, color=TEAL,
        align=PP_ALIGN.CENTER)

    notes(sl, "Thank you for watching. Full materials in micro-lesson/course-4/.")
    return sl


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
BUILDERS = [
    slide_01, slide_02, slide_03, slide_04, slide_05, slide_06,
    slide_07, slide_08, slide_09, slide_10, slide_11,
]

def build():
    prs = new_prs()
    for fn in BUILDERS:
        fn(prs)
    out = Path(__file__).parent / "sqlite-search-cache.pptx"
    prs.save(str(out))
    print(f"Saved:  {out}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Size:   {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")


if __name__ == "__main__":
    build()
